#!/usr/bin/env python3

import argparse
import json
import mimetypes
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Tuple
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MEDIA_BLOCK_TITLE = "对应图片路径："
WEB_RENDERABLE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp"}
RELATIONSHIP_NAMESPACE = {"pr": "http://schemas.openxmlformats.org/package/2006/relationships"}


@dataclass
class PictureShapeRef:
    slide_number: int
    rel_id: str
    left: float
    top: float
    width: float
    height: float
    area: float


def emu_to_inches(value: int) -> float:
    return round(value / 914400.0, 2)


def repo_relative_path(path: Path, repo_root: Path) -> str:
    resolved = path.resolve()
    try:
        return resolved.relative_to(repo_root.resolve()).as_posix()
    except ValueError:
        return resolved.as_posix()


def repo_local_path(value: str, repo_root: Path) -> Path:
    path = Path(value)
    if path.is_absolute():
        return path
    return repo_root / path


def walk_picture_shapes(shapes: Iterable, slide_number: int, refs: List[PictureShapeRef]) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_picture_shapes(shape.shapes, slide_number, refs)
            continue
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        try:
            rel_id = shape._element.blipFill.blip.rEmbed
        except Exception:
            continue

        width = emu_to_inches(shape.width)
        height = emu_to_inches(shape.height)
        refs.append(
            PictureShapeRef(
                slide_number=slide_number,
                rel_id=rel_id,
                left=emu_to_inches(shape.left),
                top=emu_to_inches(shape.top),
                width=width,
                height=height,
                area=round(width * height, 2),
            )
        )


def collect_picture_shapes(pptx_path: Path) -> Dict[Tuple[int, str], List[PictureShapeRef]]:
    prs = Presentation(str(pptx_path))
    by_slide_rel: Dict[Tuple[int, str], List[PictureShapeRef]] = {}

    for slide_number, slide in enumerate(prs.slides, start=1):
        refs: List[PictureShapeRef] = []
        walk_picture_shapes(slide.shapes, slide_number, refs)
        for ref in refs:
            by_slide_rel.setdefault((ref.slide_number, ref.rel_id), []).append(ref)

    return by_slide_rel


def dedupe_media_refs(refs: List[dict]) -> List[dict]:
    seen = set()
    output: List[dict] = []
    for ref in refs:
        key = ref["path"]
        if key in seen:
            continue
        seen.add(key)
        output.append(ref)
    return output


def choose_primary_media_refs(refs: List[dict]) -> List[dict]:
    if not refs:
        return []

    huge = [ref for ref in refs if ref["area"] >= 6]
    large = [ref for ref in refs if ref["area"] >= 2]
    medium = [ref for ref in refs if ref["area"] >= 1]

    if huge:
        return huge[:3]
    if large:
        return large[:3]
    if len(refs) <= 2 and medium:
        return medium[:2]
    if medium and refs[0]["area"] >= 1.25:
        return [medium[0]]
    return []


def strip_existing_media_block(text: str) -> str:
    marker = f"\n\n{MEDIA_BLOCK_TITLE}\n"
    if marker not in text:
        return text.rstrip()
    return text.split(marker, 1)[0].rstrip()


def append_media_block(text: str, primary_refs: List[dict]) -> str:
    base = strip_existing_media_block(text)
    if not primary_refs:
        return base

    lines = [MEDIA_BLOCK_TITLE]
    for ref in primary_refs:
        lines.append(f"- {ref['path']}")
    return f"{base}\n\n" + "\n".join(lines)


def extract_media_from_archive(pptx_path: Path, media_root: Path) -> List[str]:
    media_root.mkdir(parents=True, exist_ok=True)

    extracted_assets: List[str] = []
    with ZipFile(pptx_path) as archive:
        for member in sorted(archive.namelist()):
            if not member.startswith("ppt/media/") or member.endswith("/"):
                continue

            asset_file = Path(member).name
            target_path = media_root / asset_file
            target_path.write_bytes(archive.read(member))
            extracted_assets.append(asset_file)

    expected = set(extracted_assets)
    for existing in media_root.iterdir():
        if not existing.is_file():
            continue
        if existing.name == "index.json":
            continue
        if existing.name not in expected:
            existing.unlink()

    return extracted_assets


def build_media_index(
    pptx_path: Path,
    media_root: Path,
    source_path: Path,
    repo_root: Path,
) -> dict:
    slide_rel_pattern = re.compile(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$")
    items: List[dict] = []

    with ZipFile(pptx_path) as archive:
        for member in sorted(archive.namelist()):
            match = slide_rel_pattern.fullmatch(member)
            if not match:
                continue

            slide_number = int(match.group(1))
            root = ET.fromstring(archive.read(member))
            for rel in root.findall("pr:Relationship", RELATIONSHIP_NAMESPACE):
                relationship_type = rel.attrib.get("Type", "")
                short_type = relationship_type.rsplit("/", 1)[-1]
                if short_type not in {"image", "media", "hdphoto"}:
                    continue

                asset_file = Path(rel.attrib.get("Target", "")).name
                asset_path = media_root / asset_file
                if not asset_path.exists():
                    continue

                ext = asset_path.suffix.lower()
                items.append(
                    {
                        "id": f"slide-{slide_number}-{rel.attrib.get('Id', '')}-{asset_file}",
                        "slideNumber": slide_number,
                        "relId": rel.attrib.get("Id", ""),
                        "relationshipType": short_type,
                        "assetFile": asset_file,
                        "path": repo_relative_path(asset_path, repo_root),
                        "sourcePath": repo_relative_path(source_path, repo_root),
                        "webRenderable": ext in WEB_RENDERABLE_EXTENSIONS,
                        "mimeType": mimetypes.guess_type(asset_file)[0],
                    }
                )

    items.sort(key=lambda item: (item["slideNumber"], item["assetFile"], item["relId"]))
    asset_count = sum(
        1
        for asset_path in media_root.iterdir()
        if asset_path.is_file() and asset_path.name != "index.json"
    )

    return {
        "sourcePath": repo_relative_path(source_path, repo_root),
        "rootPath": repo_relative_path(media_root, repo_root),
        "indexPath": repo_relative_path(media_root / "index.json", repo_root),
        "assetCount": asset_count,
        "referenceCount": len(items),
        "webRenderableReferenceCount": sum(1 for item in items if item["webRenderable"]),
        "relationshipTypeCounts": {
            key: sum(1 for item in items if item["relationshipType"] == key)
            for key in sorted({item["relationshipType"] for item in items})
        },
        "items": items,
    }


def resolve_inputs(
    normalized: dict,
    repo_root: Path,
    source_archive_dir: str | None,
    pptx_file: str | None,
) -> Tuple[Path, Path]:
    stored_path = normalized.get("importedSource", {}).get("storedPath")
    default_pptx_path = repo_local_path(stored_path, repo_root) if stored_path else None

    resolved_pptx_path = repo_local_path(pptx_file, repo_root) if pptx_file else default_pptx_path
    if resolved_pptx_path is None or not resolved_pptx_path.exists():
        raise FileNotFoundError("Unable to locate archived PPTX file. Pass --pptx-file to override.")

    resolved_archive_dir = (
        repo_local_path(source_archive_dir, repo_root)
        if source_archive_dir
        else resolved_pptx_path.parent
    )
    resolved_archive_dir.mkdir(parents=True, exist_ok=True)

    return resolved_archive_dir, resolved_pptx_path


def annotate_sections_with_media(
    normalized: dict,
    media_index: dict,
    picture_shapes: Dict[Tuple[int, str], List[PictureShapeRef]],
) -> None:
    index_by_slide_rel = {
        (item["slideNumber"], item["relId"]): item
        for item in media_index.get("items", [])
        if item.get("webRenderable")
    }

    for section in normalized.get("sections", []):
        slide_number = section.get("order")
        refs: List[dict] = []

        for (shape_slide_number, rel_id), shapes in picture_shapes.items():
            if shape_slide_number != slide_number:
                continue

            media_item = index_by_slide_rel.get((shape_slide_number, rel_id))
            if not media_item:
                continue

            for shape in shapes:
                refs.append(
                    {
                        "path": media_item["path"],
                        "assetFile": media_item["assetFile"],
                        "relId": rel_id,
                        "slideNumber": shape.slide_number,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "area": shape.area,
                    }
                )

        refs.sort(key=lambda item: (-item["area"], item["top"], item["left"], item["assetFile"]))
        refs = dedupe_media_refs(refs)

        primary_paths = {ref["path"] for ref in choose_primary_media_refs(refs)}
        for ref in refs:
            ref["role"] = "primary" if ref["path"] in primary_paths else "supporting"

        section["mediaRefs"] = refs
        primary_refs = [ref for ref in refs if ref["role"] == "primary"]
        section["content"] = append_media_block(section.get("content", ""), primary_refs)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Extract PPTX media into the source archive, build a media index, and annotate normalized sections."
    )
    parser.add_argument("normalized_json", help="Path to content/normalized/<source-id>.json")
    parser.add_argument(
        "--source-archive-dir",
        help="Optional override for content/sources/<kind>/<year>/<archive-key>/",
    )
    parser.add_argument("--pptx-file", help="Optional override for the archived PPTX file path")
    args = parser.parse_args()

    repo_root = Path.cwd().resolve()
    normalized_path = repo_local_path(args.normalized_json, repo_root)
    normalized = json.loads(normalized_path.read_text(encoding="utf-8"))

    source_archive_dir, pptx_path = resolve_inputs(
        normalized,
        repo_root,
        args.source_archive_dir,
        args.pptx_file,
    )
    media_root = source_archive_dir / "media"
    extract_media_from_archive(pptx_path, media_root)

    media_index = build_media_index(
        pptx_path=pptx_path,
        media_root=media_root,
        source_path=pptx_path,
        repo_root=repo_root,
    )
    (media_root / "index.json").write_text(
        json.dumps(media_index, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )

    picture_shapes = collect_picture_shapes(pptx_path)
    annotate_sections_with_media(normalized, media_index, picture_shapes)
    normalized["media"] = {
        key: media_index[key]
        for key in [
            "rootPath",
            "indexPath",
            "assetCount",
            "referenceCount",
            "webRenderableReferenceCount",
            "relationshipTypeCounts",
        ]
    }
    normalized_path.write_text(json.dumps(normalized, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    total_sections = len(normalized.get("sections", []))
    sections_with_refs = sum(1 for section in normalized.get("sections", []) if section.get("mediaRefs"))
    sections_with_primary_paths = sum(
        1 for section in normalized.get("sections", []) if f"\n\n{MEDIA_BLOCK_TITLE}\n" in section.get("content", "")
    )

    print(
        json.dumps(
            {
                "normalizedPath": repo_relative_path(normalized_path, repo_root),
                "pptxPath": repo_relative_path(pptx_path, repo_root),
                "mediaRoot": media_index["rootPath"],
                "mediaIndexPath": media_index["indexPath"],
                "assetCount": media_index["assetCount"],
                "referenceCount": media_index["referenceCount"],
                "webRenderableReferenceCount": media_index["webRenderableReferenceCount"],
                "totalSections": total_sections,
                "sectionsWithMediaRefs": sections_with_refs,
                "sectionsWithInlineMediaPaths": sections_with_primary_paths,
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
